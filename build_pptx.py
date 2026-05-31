"""PowerPoint 生成モジュール（操作マニュアル・導入説明資料）"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット ──────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x56, 0xDB)
NAVY_D  = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_L  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT  = RGBColor(0xDB, 0xEA, 0xFE)
ACCENT2 = RGBColor(0xEF, 0xF6, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x05, 0x96, 0x69)
GREEN_D = RGBColor(0x06, 0x5F, 0x46)
GREEN_L = RGBColor(0xD1, 0xFA, 0xE5)
ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_D= RGBColor(0x92, 0x40, 0x0E)
ORANGE_L= RGBColor(0xFE, 0xF3, 0xC7)
GRAY    = RGBColor(0xF1, 0xF5, 0xF9)
GRAY2   = RGBColor(0xE2, 0xE8, 0xF0)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
MUTED_L = RGBColor(0x94, 0xA3, 0xB8)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_D= RGBColor(0x5B, 0x21, 0xB6)
PURPLE_L= RGBColor(0xED, 0xE9, 0xFE)
TEAL    = RGBColor(0x0F, 0x76, 0x6E)
TEAL_L  = RGBColor(0xCC, 0xFB, 0xF1)
PINK    = RGBColor(0xBE, 0x18, 0x5D)
PINK_L  = RGBColor(0xFF, 0xF0, 0xF7)
RED_L   = RGBColor(0xFE, 0xE2, 0xE2)
RED     = RGBColor(0xEF, 0x44, 0x44)

# ── スライドサイズ ──────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
HH = Inches(0.78)          # header height
FH = Inches(0.32)          # footer height
ML = Inches(0.55)          # margin left/right
CY = HH + Inches(0.28)     # content start Y
CH = H - CY - FH - Inches(0.18)   # content height
CW = W - ML * 2            # content width


# ── 低レベルヘルパー ─────────────────────────────────────────
def _prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _box(slide, x, y, w, h, rgb, line_rgb=None, line_w=0.5, shape_type=1):
    s = slide.shapes.add_shape(shape_type, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def _rbox(slide, x, y, w, h, rgb, line_rgb=None, line_w=0.5):
    """角丸長方形"""
    return _box(slide, x, y, w, h, rgb, line_rgb, line_w, shape_type=5)

def _oval(slide, x, y, w, h, rgb, line_rgb=None):
    return _box(slide, x, y, w, h, rgb, line_rgb, shape_type=9)

def _pentagon(slide, x, y, w, h, rgb):
    return _box(slide, x, y, w, h, rgb, shape_type=51)

def _chevron(slide, x, y, w, h, rgb):
    return _box(slide, x, y, w, h, rgb, shape_type=52)

def _set_tf(shape, lines, size=11, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, ml=Inches(0.18), mt=Inches(0.1),
            mr=Inches(0.1), mb=Inches(0.05), italic=False, line_sp=None):
    tf = shape.text_frame
    tf.word_wrap     = True
    tf.margin_left   = int(ml)
    tf.margin_top    = int(mt)
    tf.margin_right  = int(mr)
    tf.margin_bottom = int(mb)
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            from pptx.util import Pt as _Pt
            from pptx.oxml.ns import qn
            from lxml import etree
            lnSpc = etree.SubElement(p._p, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_sp * 100)))
        r = p.add_run()
        r.text        = line
        r.font.size   = Pt(size)
        r.font.bold   = bold
        r.font.italic = italic
        r.font.name   = "Meiryo"
        r.font.color.rgb = color

def _tb(slide, x, y, w, h, text, size=11, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    _set_tf(box, text, size=size, bold=bold, color=color, align=align,
            italic=italic, ml=Inches(0.05), mt=Inches(0.03))
    return box

def _chrome(slide, title, n, total=10):
    """ヘッダー＋フッター共通フレーム"""
    # 背景
    _box(slide, 0, 0, W, H, GRAY)
    # ヘッダー
    _box(slide, 0, 0, W, HH, NAVY_D)
    _box(slide, 0, 0, Inches(0.1), HH, NAVY_L)
    hdr = _box(slide, Inches(0.18), 0, W - Inches(0.18), HH, NAVY_D)
    _set_tf(hdr, title, size=19, bold=True, color=WHITE,
            ml=Inches(0.15), mt=Inches(0.16))
    num = _box(slide, W - Inches(1.5), Inches(0.2), Inches(1.3), Inches(0.38), NAVY_D)
    _set_tf(num, f"{n}  /  {total}", size=10, color=MUTED_L,
            align=PP_ALIGN.RIGHT, ml=Inches(0.05), mt=Inches(0.06))
    # コンテンツ背景 (白カード)
    _box(slide, ML - Inches(0.1), CY - Inches(0.1),
         CW + Inches(0.2), CH + Inches(0.25), WHITE,
         line_rgb=GRAY2, line_w=0.5)
    # フッター
    _box(slide, 0, H - FH, W, FH, NAVY_D)

def _section_badge(slide, x, y, w, h, text, bg=NAVY, tc=WHITE, size=13):
    s = _rbox(slide, x, y, w, h, bg)
    _set_tf(s, text, size=size, bold=True, color=tc,
            align=PP_ALIGN.CENTER, ml=Inches(0.1), mt=Inches(0.08))
    return s

def _card(slide, x, y, w, h, bg=WHITE, border=GRAY2):
    s = _rbox(slide, x, y, w, h, bg, line_rgb=border, line_w=0.75)
    return s

def _accent_card(slide, x, y, w, h, accent_color, bg=WHITE):
    """左側にカラーバーを持つカード"""
    _rbox(slide, x, y, w, h, bg, line_rgb=GRAY2, line_w=0.5)
    _box(slide, x, y, Inches(0.06), h, accent_color)
    return slide.shapes[-2]

def _num_badge(slide, x, y, size_px, num_str, bg, tc=WHITE):
    o = _oval(slide, x, y, size_px, size_px, bg)
    _set_tf(o, num_str, size=11, bold=True, color=tc,
            align=PP_ALIGN.CENTER, ml=Inches(0.02), mt=Inches(0.04))
    return o

def _step_flow(slide, steps, x, y, w, h):
    """ステップフロー（pentagon + chevron）"""
    n = len(steps)
    sw = (w - Inches(0.05) * (n - 1)) / n
    for i, (label, title, color, optional) in enumerate(steps):
        bx = x + (sw + Inches(0.05)) * i
        shape_fn = _pentagon if i == 0 else _chevron
        s = shape_fn(slide, bx, y, sw, h, color)
        _set_tf(s, [label, title], size=8, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.08), mt=Inches(0.06))
    return sw


# ════════════════════════════════════════════════════
#  操作マニュアル
# ════════════════════════════════════════════════════
def _build_manual_pptx() -> io.BytesIO:
    prs = _prs()
    TOTAL = 10

    # ── Slide 1 : タイトル ─────────────────────────────────
    sl = _slide(prs)
    _box(sl, 0, 0, W, H, NAVY_D)
    # 装飾：右下の大きな円
    _oval(sl, W - Inches(4.5), H - Inches(4.5), Inches(6), Inches(6),
          RGBColor(0x1E, 0x40, 0xAF))
    _oval(sl, W - Inches(3.2), H - Inches(3.2), Inches(4), Inches(4),
          RGBColor(0x1D, 0x4E, 0xD8))
    # 装飾：左上の小さな円
    _oval(sl, -Inches(1.0), -Inches(1.0), Inches(3), Inches(3),
          RGBColor(0x1E, 0x40, 0xAF))
    # タイトルエリア
    _box(sl, 0, Inches(2.2), Inches(0.15), Inches(1.8), NAVY_L)
    title_b = _box(sl, Inches(0.3), Inches(2.2), W - Inches(4.5), Inches(1.8), NAVY_D)
    _set_tf(title_b, "日報自動入力アプリ\n操作マニュアル",
            size=34, bold=True, color=WHITE, ml=Inches(0.3), mt=Inches(0.2))
    sub_b = _box(sl, Inches(0.3), Inches(4.2), W - Inches(4.5), Inches(0.55), NAVY_D)
    _set_tf(sub_b, "このマニュアルを読めば、月初の作業が 5分で完了します",
            size=13, color=RGBColor(0xBF, 0xD7, 0xFF), ml=Inches(0.3), mt=Inches(0.06))
    # ガイド内容
    for i, line in enumerate([
        "◆  アプリの使い方（STEP 1〜7）の詳しい手順",
        "◆  自動入力ルールの早見表",
        "◆  よくある疑問と回答",
    ]):
        _tb(sl, Inches(0.55), Inches(5.1) + Inches(0.42) * i,
            W - Inches(4.5), Inches(0.38), line,
            size=11, color=RGBColor(0x93, 0xC5, 0xFD))
    # フッター
    _box(sl, 0, H - Inches(0.32), W, Inches(0.32), NAVY_D)
    _tb(sl, W - Inches(1.5), H - Inches(0.28), Inches(1.3), Inches(0.25),
        "1  /  10", size=9, color=MUTED_L, align=PP_ALIGN.RIGHT)

    # ── Slide 2 : このアプリについて ───────────────────────
    sl = _slide(prs)
    _chrome(sl, "このアプリについて", 2, TOTAL)
    cy = CY + Inches(0.05)
    fw = (CW - Inches(0.3)) / 4

    features = [
        (NAVY,   "📅", "月初 5分以内",    "で操作を完了できます",         ACCENT2),
        (GREEN,  "📊", "1ヶ月分すべて",   "のExcel行を自動で埋めます",     GREEN_L),
        (TEAL,   "🗓", "土日・祝日・有給", "を自動で判定・記入します",      TEAL_L),
        (PURPLE, "💾", "設定は自動保存",   "翌月は変更箇所のみ修正OK",      PURPLE_L),
    ]
    for i, (color, icon, kw, rest, bg) in enumerate(features):
        fx = ML + (fw + Inches(0.1)) * i
        # アイコン円
        _oval(sl, fx + fw / 2 - Inches(0.38), cy, Inches(0.75), Inches(0.75), color)
        _tb(sl, fx + fw / 2 - Inches(0.38), cy, Inches(0.75), Inches(0.75),
            icon, size=16, align=PP_ALIGN.CENTER)
        # カード
        card = _rbox(sl, fx, cy + Inches(0.82), fw, Inches(1.35), bg, GRAY2)
        _set_tf(card, [kw, rest], size=11, color=DARK,
                align=PP_ALIGN.CENTER, ml=Inches(0.08), mt=Inches(0.15))
        card.text_frame.paragraphs[0].runs[0].font.bold  = True
        card.text_frame.paragraphs[0].runs[0].font.color.rgb = color
        card.text_frame.paragraphs[0].runs[0].font.size  = Pt(13)

    # 事前確認セクション
    req_y = cy + Inches(2.35)
    _section_badge(sl, ML, req_y, CW, Inches(0.38), "必要なもの（事前確認）", GREEN_D)
    checks = [
        ("✔", "ブラウザ（Chrome / Edge / Safari）",   "インストール不要"),
        ("✔", "アプリのURL・ログイン情報",              "管理者から入手してください"),
        ("✔", "会社配布の作業報告書（.xlsx）",          "毎月配布されるファイル"),
    ]
    cw3 = (CW - Inches(0.3)) / 3
    for i, (icon, main, sub) in enumerate(checks):
        cx = ML + (cw3 + Inches(0.15)) * i
        c = _rbox(sl, cx, req_y + Inches(0.45), cw3, Inches(0.9), WHITE, GRAY2)
        _set_tf(c, [f"{icon}  {main}", sub], size=10, color=DARK,
                align=PP_ALIGN.CENTER, ml=Inches(0.1), mt=Inches(0.12))
        c.text_frame.paragraphs[0].runs[0].font.bold  = True
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = GREEN_D
        c.text_frame.paragraphs[1].runs[0].font.color.rgb = MUTED
        c.text_frame.paragraphs[1].runs[0].font.size  = Pt(9)

    _tb(sl, ML, req_y + Inches(1.42), CW, Inches(0.28),
        "！ インターネット接続が必要です（祝日取得）　　！ ダウンロード後、内容を確認してから保存　　！ 元のファイルは上書きされません",
        size=9, color=ORANGE, italic=True)

    # ── Slide 3 : 操作の全体フロー ─────────────────────────
    sl = _slide(prs)
    _chrome(sl, "操作の全体フロー", 3, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.3),
        "7 つのステップで月初の作業を完了させましょう　※ STEP 4・5（オレンジ色）は任意です",
        size=10, color=MUTED, italic=True)

    flow_steps = [
        ("STEP 1", "アクセス\nログイン",    NAVY,    False),
        ("STEP 2", "対象月\nを確認",        GREEN_D, False),
        ("STEP 3", "勤務時間\n設定",         ORANGE_D,False),
        ("STEP 4", "例外日\n設定（任意）",   ORANGE,  True),
        ("STEP 5", "有給日\n選択（任意）",   PINK,    True),
        ("STEP 6", "Excel\nファイル選択",    TEAL,    False),
        ("STEP 7", "完了・\nDL",            NAVY,    False),
    ]
    sw = _step_flow(sl, flow_steps,
                    ML, CY + Inches(0.42), CW, Inches(1.05))

    # 説明カード
    descs = [
        (NAVY,    "必須  STEP 1〜3 ・ 6〜7",
         "毎月行う操作です。STEP 3の設定はブラウザに保存されるため、2回目以降は変更がある場合のみ修正してください。"),
        (ORANGE,  "任意  STEP 4〜5",
         "残業・早退・休日出勤がある場合は STEP 4 を。有給休暇がある場合は STEP 5 を操作してください。"),
    ]
    dw = (CW - Inches(0.3)) / 2
    dy = CY + Inches(1.65)
    for i, (color, head, body) in enumerate(descs):
        dx = ML + (dw + Inches(0.3)) * i
        _accent_card(sl, dx, dy, dw, Inches(1.5), color)
        hb = _box(sl, dx + Inches(0.08), dy, dw - Inches(0.08), Inches(0.42), WHITE)
        _set_tf(hb, head, size=11, bold=True, color=color,
                ml=Inches(0.15), mt=Inches(0.1))
        bb = _box(sl, dx + Inches(0.08), dy + Inches(0.42), dw - Inches(0.08), Inches(1.0), WHITE)
        _set_tf(bb, body, size=10, color=DARK, ml=Inches(0.15), mt=Inches(0.1))

    # ポイント注釈
    note = _rbox(sl, ML, dy + Inches(1.58), CW, Inches(0.45), ACCENT2, ACCENT)
    _set_tf(note, "設定はすべてブラウザに自動保存されます。翌月からは変更のある箇所だけ直せばOKです。",
            size=10, color=NAVY, ml=Inches(0.25), mt=Inches(0.12), align=PP_ALIGN.CENTER)

    # ── Slide 4 : STEP 1 ────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "STEP 1　アクセス・ログイン", 4, TOTAL)
    _num_badge(sl, ML, CY + Inches(0.05), Inches(0.58), "1", NAVY)
    _tb(sl, ML + Inches(0.68), CY + Inches(0.1), CW - Inches(0.68), Inches(0.4),
        "ブラウザでアプリのURLを開き、ログインします",
        size=13, bold=True, color=NAVY_D)
    detail_steps = [
        "①  ブラウザ（Chrome / Edge / Safari）を開きます。",
        "②  アドレスバーにアプリの URL を入力してアクセスします。",
        "③  ユーザー名・パスワードの入力画面が表示されたら、指定の情報を入力して「OK」を押します。",
        "④  アプリのメイン画面が表示されればログイン完了です。",
    ]
    colors_alt = [ACCENT2, WHITE, ACCENT2, WHITE]
    for i, (s, bg) in enumerate(zip(detail_steps, colors_alt)):
        row = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * i,
                    CW, Inches(0.68), bg, GRAY2)
        _set_tf(row, s, size=12, color=DARK, ml=Inches(0.25), mt=Inches(0.16))

    tip = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * 4 + Inches(0.1),
                CW, Inches(0.72), RGBColor(0xEF, 0xF6, 0xFF), NAVY_L, line_w=1.0)
    _set_tf(tip, ["💡  ポイント",
                  "URL とログイン情報は管理者から入手してください。ログインできない場合は Caps Lock やスペース入力ミスを確認してください。"],
            size=10, color=NAVY_D, ml=Inches(0.2), mt=Inches(0.1))
    tip.text_frame.paragraphs[0].runs[0].font.bold = True
    tip.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # ── Slide 5 : STEP 2〜3 ─────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "STEP 2〜3　対象月の確認・勤務時間の設定", 5, TOTAL)
    half = (CW - Inches(0.3)) / 2
    for col, (num_s, head_text, head_color, items, item_bg) in enumerate([
        ("2", "対象月を確認する", GREEN_D,
         ["画面上部に「年・月」が表示されます。",
          "自動で「当月」が設定されています。",
          "別の月を対象にしたい場合は数字を直接変更してください。",
          "➡ 通常は変更不要。当月が自動で設定されています。"],
         GREEN_L),
        ("3", "勤務時間を設定する", ORANGE_D,
         ["月〜金それぞれの「開始時間」「終了時間」を入力します。",
          "入力後、下部の月間スケジュールに1日〜末日の予定と合計稼働時間が表示されます。",
          "設定変更時はヘッダーに「✓ 自動保存」と表示され、ブラウザに保存されます。",
          "➡ 翌月以降は変更箇所のみ修正するだけでOKです。"],
         ORANGE_L),
    ]):
        cx = ML + (half + Inches(0.3)) * col
        _num_badge(sl, cx, CY + Inches(0.05), Inches(0.55), num_s, head_color)
        _tb(sl, cx + Inches(0.65), CY + Inches(0.1), half - Inches(0.65), Inches(0.38),
            head_text, size=13, bold=True, color=head_color)
        for i, s in enumerate(items):
            bg = item_bg if i % 2 == 0 else WHITE
            bd = RGBColor(0x6E, 0xE7, 0xB7) if col == 0 else RGBColor(0xFB, 0xBF, 0x24)
            row = _rbox(sl, cx, CY + Inches(0.6) + Inches(0.82) * i,
                        half, Inches(0.78), bg, bd if i == 3 else GRAY2)
            bold = i == 3
            c = head_color if i == 3 else DARK
            _set_tf(row, s, size=10, color=c, bold=bold,
                    ml=Inches(0.2), mt=Inches(0.16))

    # ── Slide 6 : STEP 4 ─────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "STEP 4　例外日を設定する（任意）", 6, TOTAL)
    _num_badge(sl, ML, CY + Inches(0.05), Inches(0.58), "4", PURPLE)
    _tb(sl, ML + Inches(0.68), CY + Inches(0.1), CW - Inches(0.68), Inches(0.38),
        "残業・早退・休日出勤など、通常と異なる日を個別設定します",
        size=13, bold=True, color=PURPLE_D)
    ex_steps = [
        "①  「＋ 例外日を追加」ボタンをクリックします。",
        "②  カレンダーが表示されるので、対象の日付をクリックして選択します（土日・祝日も選択可能）。",
        "③  選択した日の開始時間・終了時間を入力します。備考欄の初期値は「在宅勤務」（変更可能）。",
        "④  複数日ある場合は必要な分だけ繰り返してください。",
    ]
    for i, s in enumerate(ex_steps):
        row = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * i,
                    CW, Inches(0.68), PURPLE_L if i % 2 == 0 else WHITE, GRAY2)
        _set_tf(row, s, size=12, color=DARK, ml=Inches(0.25), mt=Inches(0.16))

    warn = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * 4 + Inches(0.1),
                 CW, Inches(0.82), RGBColor(0xFA, 0xF5, 0xFF), PURPLE, line_w=1.2)
    _set_tf(warn, [
        "⚠  重要：例外日設定は最優先で適用されます",
        "有給・祝日・土日に関わらず最優先で反映されます。休日出勤の記録にも使えます。",
        "このステップは任意です。例外日がない月はスキップしてください。",
    ], size=10, color=PURPLE_D, ml=Inches(0.2), mt=Inches(0.1))
    warn.text_frame.paragraphs[0].runs[0].font.bold = True
    warn.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # ── Slide 7 : STEP 5 ─────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "STEP 5　有給取得日を選択する（任意）", 7, TOTAL)
    _num_badge(sl, ML, CY + Inches(0.05), Inches(0.58), "5", PINK)
    _tb(sl, ML + Inches(0.68), CY + Inches(0.1), CW - Inches(0.68), Inches(0.38),
        "有給休暇を取得する日をカレンダーで選択します",
        size=13, bold=True, color=PINK)
    paid_steps = [
        "①  「有給を取得する日がある」チェックボックスをオンにします。",
        "②  当月のカレンダーが表示されます。",
        "③  有給を取得する日付をクリックして選択します（選択した日はハイライト表示されます）。",
        "④  複数日ある場合は続けてクリックしてください。もう一度クリックすると選択を解除できます。",
    ]
    for i, s in enumerate(paid_steps):
        row = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * i,
                    CW, Inches(0.68), PINK_L if i % 2 == 0 else WHITE, GRAY2)
        _set_tf(row, s, size=12, color=DARK, ml=Inches(0.25), mt=Inches(0.16))

    tip = _rbox(sl, ML, CY + Inches(0.6) + Inches(0.72) * 4 + Inches(0.1),
                CW, Inches(0.82), RGBColor(0xFF, 0xF8, 0xFB), PINK, line_w=1.0)
    _set_tf(tip, [
        "💡  ポイント",
        "土曜・日曜・祝日も選択できます（振替休日の有給取得など）。",
        "選択した日は備考欄に「私用により、休暇」と自動入力されます。このステップは任意です。",
    ], size=10, color=PINK, ml=Inches(0.2), mt=Inches(0.1))
    tip.text_frame.paragraphs[0].runs[0].font.bold = True
    tip.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # ── Slide 8 : STEP 6〜7 ──────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "STEP 6〜7　Excelファイルの選択・ダウンロード", 8, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "最後のステップです。ファイルを選んでボタンを押すだけで完了します",
        size=10, color=MUTED, italic=True)
    half8 = (CW - Inches(0.3)) / 2
    for col, (num_s, head, color, items, bg_c) in enumerate([
        ("6", "Excelファイルを選択する", TEAL,
         ["「Excelファイルを選択」エリアをクリックします。",
          "ファイル選択ダイアログが開くので、配布された作業報告書（.xlsx）を選択します。",
          "ファイル名が表示されれば選択完了です。",
          "間違えた場合は「✕」ボタンでキャンセルして選び直せます。"],
         TEAL_L),
        ("7", "完了・ダウンロード", NAVY,
         ["「入力完了・ダウンロード」ボタンをクリックします。",
          "処理が完了すると、自動入力済みのExcelファイルがダウンロードされます。",
          "ダウンロードしたファイルを開いて内容を確認してください。",
          "問題なければ所定の場所に保存して作業完了です。お疲れ様でした！"],
         ACCENT),
    ]):
        cx = ML + (half8 + Inches(0.3)) * col
        _num_badge(sl, cx, CY + Inches(0.37), Inches(0.55), num_s, color)
        _tb(sl, cx + Inches(0.65), CY + Inches(0.42), half8 - Inches(0.65), Inches(0.38),
            head, size=12, bold=True, color=color)
        for i, s in enumerate(items):
            bg = bg_c if i % 2 == 0 else WHITE
            row = _rbox(sl, cx, CY + Inches(0.95) + Inches(0.82) * i,
                        half8, Inches(0.78), bg, GRAY2)
            _set_tf(row, s, size=10, color=DARK, ml=Inches(0.18), mt=Inches(0.16))

    note = _rbox(sl, ML, H - FH - Inches(0.52), CW, Inches(0.38), ACCENT2, ACCENT)
    _set_tf(note, "元のExcelファイルは上書きされません。新しいファイルがダウンロードされます。安心して操作してください。",
            size=10, color=NAVY, align=PP_ALIGN.CENTER, ml=Inches(0.2), mt=Inches(0.08))

    # ── Slide 9 : 自動入力ルール早見表 ──────────────────────
    sl = _slide(prs)
    _chrome(sl, "自動入力ルール早見表", 9, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "日の種類に応じて以下のルールで Excel へ自動入力されます",
        size=10, color=MUTED, italic=True)

    col_w = [Inches(2.6), Inches(4.7), Inches(4.5)]
    headers9 = ["日の種類", "開始・終了・休憩時間", "備考欄の記入内容"]
    hx = ML
    hy9 = CY + Inches(0.4)
    for w, h in zip(col_w, headers9):
        cell = _box(sl, hx, hy9, w, Inches(0.45), NAVY_D)
        _set_tf(cell, h, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.1))
        hx += w

    table_data = [
        ("出勤日（平日）",
         "曜日別設定の開始・終了時間\n休憩時間：1時間00分",
         "在宅勤務",
         ACCENT2, NAVY_L),
        ("例外日\n（残業・早退・休日出勤）",
         "例外日設定の開始・終了時間\n休憩時間：1時間00分",
         "在宅勤務\n（例外設定が最優先）",
         ORANGE_L, ORANGE),
        ("有給取得日",
         "空欄",
         "私用により、休暇",
         PINK_L, PINK),
        ("祝日",
         "空欄",
         "祝日",
         RED_L, RED),
        ("土日",
         "空欄",
         "空欄",
         GRAY, MUTED_L),
    ]
    row_h9 = (CH - Inches(0.8)) / len(table_data)
    for ri, (kind, times, note_text, bg, accent_c) in enumerate(table_data):
        ry = hy9 + Inches(0.45) + row_h9 * ri
        rx = ML
        for ci, (val, w) in enumerate([(kind, col_w[0]), (times, col_w[1]), (note_text, col_w[2])]):
            cell = _rbox(sl, rx, ry, w, row_h9, bg, GRAY2)
            if ci == 0:
                _box(sl, rx, ry, Inches(0.06), row_h9, accent_c)
            lines = val.split('\n')
            _set_tf(cell, lines, size=11, color=DARK,
                    align=PP_ALIGN.CENTER, ml=Inches(0.12), mt=Inches(0.1))
            rx += w

    notes9_y = hy9 + Inches(0.45) + row_h9 * 5 + Inches(0.05)
    for n in [
        "※ 祝日情報は外部API（jpholiday）から自動取得。インターネット未接続時は空欄になる場合があります。",
        "※ 例外日設定は有給・祝日・土日に関わらず最優先で適用されます（休日出勤の記録にも使用可能）。",
    ]:
        _tb(sl, ML, notes9_y, CW, Inches(0.26), n, size=9, color=MUTED, italic=True)
        notes9_y += Inches(0.28)

    # ── Slide 10 : Q&A ──────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "よくある疑問（Q&A）", 10, TOTAL)
    qas = [
        ("Q1", "設定は毎月入力し直す必要がありますか？",
         "不要です。曜日別の勤務時間・設定はブラウザに自動保存されます。次回起動時に前回の設定が自動で読み込まれます。"),
        ("Q2", "休日出勤はどうやって入力しますか？",
         "「＋ 例外日を追加」から対象日（土日・祝日も可）を選択し、出勤時間を入力します。例外日設定が最優先で反映されます。"),
        ("Q3", "有給を土日・祝日に取得する場合は？",
         "STEP 5のカレンダーで土日・祝日も選択できます。選択した日の備考欄に「私用により、休暇」が記入されます。"),
        ("Q4", "ファイルを間違えて選択した場合は？",
         "ファイル名の右側の「✕」ボタンで選択をキャンセルできます。再度クリックして正しいファイルを選び直してください。"),
        ("Q5", "Excelの列の位置が変わっても使えますか？",
         "はい。アプリがヘッダー文字列を自動検索して対象列を特定するため、列の位置が変わっても正常に動作します。"),
        ("Q6", "ダウンロードしたファイルはどこに保存される？",
         "ブラウザの設定に従い、通常はダウンロードフォルダに保存されます。確認後、所定の場所へ移動してください。"),
    ]
    qa_h10 = CH / len(qas)
    for i, (qn_s, q, a) in enumerate(qas):
        qy = CY + qa_h10 * i
        badge = _oval(sl, ML, qy + qa_h10 * 0.05, Inches(0.52), Inches(0.52), PURPLE)
        _set_tf(badge, qn_s, size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.03), mt=Inches(0.06))
        qrow = _rbox(sl, ML + Inches(0.62), qy + qa_h10 * 0.04,
                     CW - Inches(0.62), qa_h10 * 0.42, PURPLE_L, GRAY2)
        _set_tf(qrow, q, size=11, bold=True, color=PURPLE_D,
                ml=Inches(0.2), mt=Inches(0.08))
        arow = _rbox(sl, ML + Inches(0.62), qy + qa_h10 * 0.46,
                     CW - Inches(0.62), qa_h10 * 0.5, WHITE, GRAY2)
        _set_tf(arow, a, size=10, color=DARK, ml=Inches(0.2), mt=Inches(0.08))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════════════
#  導入説明資料（_build_report_pptx は下に続く）
# ════════════════════════════════════════════════════
#  導入説明資料
# ════════════════════════════════════════════════════
def _build_report_pptx() -> io.BytesIO:
    from datetime import date as _date
    TODAY = f"{_date.today().year}年{_date.today().month}月"
    prs = _prs()
    TOTAL = 10

    # ── Slide 1 : タイトル ─────────────────────────────────
    sl = _slide(prs)
    _box(sl, 0, 0, W, H, NAVY_D)
    _oval(sl, W - Inches(5.0), H - Inches(5.0), Inches(7), Inches(7),
          RGBColor(0x1E, 0x40, 0xAF))
    _oval(sl, W - Inches(3.5), H - Inches(3.5), Inches(5), Inches(5),
          RGBColor(0x1D, 0x4E, 0xD8))
    _oval(sl, -Inches(1.5), -Inches(1.0), Inches(4), Inches(4),
          RGBColor(0x1E, 0x40, 0xAF))
    _box(sl, 0, Inches(1.8), Inches(0.18), Inches(2.5), NAVY_L)
    t = _box(sl, Inches(0.38), Inches(1.8), W - Inches(4.5), Inches(1.4), NAVY_D)
    _set_tf(t, "日報自動入力アプリ", size=36, bold=True, color=WHITE,
            ml=Inches(0.25), mt=Inches(0.2))
    sub = _box(sl, Inches(0.38), Inches(3.3), W - Inches(4.5), Inches(0.65), NAVY_D)
    _set_tf(sub, "導入説明資料", size=22, bold=True,
            color=RGBColor(0xBF, 0xD7, 0xFF), ml=Inches(0.25), mt=Inches(0.08))
    cap = _box(sl, Inches(0.38), Inches(4.1), W - Inches(4.5), Inches(0.45), NAVY_D)
    _set_tf(cap, "― 月次作業報告書の自動化による業務効率化 ―",
            size=12, color=RGBColor(0x93, 0xC5, 0xFD),
            italic=True, ml=Inches(0.25), mt=Inches(0.05))
    for i, (label, val) in enumerate([("作成日", TODAY), ("作成者", ""), ("提出先", "")]):
        iy = Inches(5.1) + Inches(0.52) * i
        lb = _rbox(sl, Inches(3.5), iy, Inches(1.6), Inches(0.42), NAVY)
        _set_tf(lb, label, size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.1), mt=Inches(0.09))
        vb = _rbox(sl, Inches(5.2), iy, Inches(4.8), Inches(0.42),
                   RGBColor(0x1E, 0x3A, 0x8A))
        _set_tf(vb, val, size=11, color=WHITE, ml=Inches(0.15), mt=Inches(0.09))
    _box(sl, 0, H - Inches(0.32), W, Inches(0.32), NAVY_D)
    _tb(sl, W - Inches(1.5), H - Inches(0.28), Inches(1.3), Inches(0.25),
        "1  /  10", size=9, color=MUTED_L, align=PP_ALIGN.RIGHT)

    # ── Slide 2 : 目次 ────────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "目次　CONTENTS", 2, TOTAL)
    items = [
        ("1", "現状の課題",          NAVY),
        ("2", "アプリ概要・解決策",   GREEN_D),
        ("3", "期待される効果",       TEAL),
        ("4", "動作環境・対象ユーザー", PURPLE),
        ("5", "操作手順（全体フロー）", ORANGE_D),
        ("6", "自動入力ルール",       PINK),
        ("7", "Q&A",                 NAVY),
        ("8", "まとめ・導入のお願い", GREEN_D),
    ]
    iw = (CW - Inches(0.3)) / 2
    ih = (CH - Inches(0.1)) / 4
    for idx, (num_s, title, color) in enumerate(items):
        col = idx % 2
        row = idx // 2
        ix = ML + col * (iw + Inches(0.3))
        iy = CY + Inches(0.05) + ih * row
        card = _rbox(sl, ix, iy, iw, ih - Inches(0.1), GRAY, GRAY2)
        _box(sl, ix, iy, Inches(0.06), ih - Inches(0.1), color)
        badge = _oval(sl, ix + Inches(0.18), iy + (ih - Inches(0.1)) / 2 - Inches(0.3),
                      Inches(0.6), Inches(0.6), color)
        _set_tf(badge, num_s, size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.1))
        _tb(sl, ix + Inches(0.88), iy + (ih - Inches(0.1)) / 2 - Inches(0.2),
            iw - Inches(0.95), Inches(0.42),
            title, size=13, bold=True, color=color)

    # ── Slide 3 : 現状の課題 ─────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "1.  現状の課題", 3, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "毎月発生している手作業の実態と、本アプリによる解決", size=10, color=MUTED, italic=True)
    half3 = (CW - Inches(0.4)) / 2

    for col, (head, head_color, items3, item_bg) in enumerate([
        ("✗  現状の課題", RED,
         ["毎月末〜月初に1ヶ月分の勤務時間をExcelへ手作業で入力",
          "開始・終了・休憩・備考を全日付分繰り返す単調な入力作業",
          "手入力ミス・記入漏れのリスクが常に存在する",
          "月あたり約15〜30分の入力工数が継続的に発生"],
         RED_L),
        ("✔  導入後の姿", GREEN,
         ["月初に数分の確認だけで1ヶ月分の作業がすべて完結",
          "ボタン1つで全日程を自動入力・Excelファイル生成",
          "ルールベース処理でミス・漏れを根本から排除",
          "翌月以降は設定の確認・変更のみで対応可能"],
         GREEN_L),
    ]):
        cx = ML + (half3 + Inches(0.4)) * col
        h = _rbox(sl, cx, CY + Inches(0.38), half3, Inches(0.48), head_color)
        _set_tf(h, head, size=14, bold=True, color=WHITE,
                ml=Inches(0.2), mt=Inches(0.1))
        for i, s in enumerate(items3):
            bg = item_bg if i % 2 == 0 else WHITE
            row = _rbox(sl, cx, CY + Inches(0.38) + Inches(0.48) + Inches(0.75) * i,
                        half3, Inches(0.72), bg, GRAY2)
            _set_tf(row, s, size=11, color=DARK, ml=Inches(0.2), mt=Inches(0.15))

    arrow = _tb(sl, ML + half3 + Inches(0.08), CY + Inches(1.5),
                Inches(0.25), Inches(0.5), "▶▶", size=14, bold=True,
                color=ORANGE, align=PP_ALIGN.CENTER)

    # ── Slide 4 : アプリ概要 ─────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "2.  アプリ概要・解決策", 4, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "ブラウザだけで完結する Excel 自動入力 Web アプリ", size=12, bold=True, color=NAVY)
    _tb(sl, ML, CY + Inches(0.35), CW, Inches(0.3),
        "会社から毎月配布されるExcel作業報告書への記入を自動化するWebアプリケーションです。インストール不要・ブラウザのみで動作。",
        size=10, color=MUTED, italic=True)
    features4 = [
        (NAVY,   "📅", "曜日別 勤務時間設定",
         "月〜金それぞれの開始・終了時間を設定。設定変更時はヘッダーに「自動保存」を表示。リセットボタンで初期値に戻せます。"),
        (GREEN_D,"🗓", "例外日・有給 個別設定",
         "残業・早退・休日出勤・有給取得日を個別設定。例外日は祝日・土日を上書きするため休日出勤にも対応します。"),
        (TEAL,   "📊", "月間スケジュール プレビュー",
         "1日〜末日の予定を自動計算して一覧表示。就業時間合計もリアルタイムで確認できます。"),
        (PURPLE, "💾", "Excelへ 自動書込み",
         "配布ファイルを選ぶだけ。自動入力済みファイルをダウンロード。元ファイルは上書きされません。"),
    ]
    fw4 = (CW - Inches(0.3)) / 4
    fy4 = CY + Inches(0.72)
    for i, (color, icon, title4, desc4) in enumerate(features4):
        fx = ML + (fw4 + Inches(0.1)) * i
        icon_card = _rbox(sl, fx, fy4, fw4, Inches(0.58), color)
        _set_tf(icon_card, f"{icon}  {title4}", size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.14))
        desc_card = _rbox(sl, fx, fy4 + Inches(0.58), fw4, Inches(1.9), WHITE, GRAY2)
        _set_tf(desc_card, desc4, size=10, color=DARK,
                ml=Inches(0.12), mt=Inches(0.12))

    # ── Slide 5 : 期待される効果 ─────────────────────────
    sl = _slide(prs)
    _chrome(sl, "3.  期待される効果", 5, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "導入によって得られる定量・定性効果", size=10, color=MUTED, italic=True)
    effects = [
        (NAVY,   "⏱",  "工数削減",     "約15〜30分/月",  "月次の手入力をほぼゼロに。\n年間で3〜6時間分の業務時間を創出します。"),
        (GREEN,  "✔",  "入力ミス防止",  "精度 100%",     "手動入力によるミス・漏れを排除。\n正確な作業報告書を自動生成します。"),
        (ORANGE, "🖱",  "操作の簡便性",  "月初 5分以内",  "ブラウザのみ、インストール不要。\nプログラミング知識不要で完結します。"),
        (PURPLE, "💾",  "設定の引き継ぎ", "翌月は確認のみ", "設定がブラウザに自動保存。\n翌月以降は変更箇所のみ修正すればOKです。"),
    ]
    ew = (CW - Inches(0.3)) / 4
    ey = CY + Inches(0.38)
    for i, (color, icon, title5, metric, desc5) in enumerate(effects):
        ex5 = ML + (ew + Inches(0.1)) * i
        icon_b = _oval(sl, ex5 + ew / 2 - Inches(0.4), ey, Inches(0.8), Inches(0.8), color)
        _set_tf(icon_b, icon, size=18, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.1))
        met_b = _rbox(sl, ex5, ey + Inches(0.88), ew, Inches(0.62), color)
        _set_tf(met_b, metric, size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.09))
        tit_b = _rbox(sl, ex5, ey + Inches(1.5), ew, Inches(0.38), GRAY)
        _set_tf(tit_b, title5, size=11, bold=True, color=color,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.07))
        desc_b = _rbox(sl, ex5, ey + Inches(1.88), ew, Inches(1.3), WHITE, GRAY2)
        lines5 = desc5.split('\n')
        _set_tf(desc_b, lines5, size=10, color=DARK,
                align=PP_ALIGN.CENTER, ml=Inches(0.08), mt=Inches(0.1))

    summary5 = _rbox(sl, ML, ey + Inches(3.28), CW, Inches(0.52), NAVY_D)
    _set_tf(summary5,
            "月初5分以内で1ヶ月分が自動完成。手入力ゼロ・ミスゼロを実現します。",
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            ml=Inches(0.2), mt=Inches(0.12))

    # ── Slide 6 : 動作環境・対象ユーザー ─────────────────
    sl = _slide(prs)
    _chrome(sl, "4.  動作環境・対象ユーザー", 6, TOTAL)
    env_rows = [
        (NAVY,    "対象ユーザー", "Excel形式（.xlsx）の作業報告書を毎月提出している社員"),
        (GREEN_D, "対象ファイル", "会社から毎月配布される Excel 作業報告書（.xlsx）"),
        (TEAL,    "動作環境",    "ブラウザ（Chrome / Edge / Safari など）　※ インストール不要"),
        (PURPLE,  "アクセス方法", "社内ネットワーク上のURL、またはローカル環境からアクセス"),
        (ORANGE_D,"認証方式",    "ユーザー名・パスワードによる Basic 認証（不正アクセス防止）"),
        (PINK,    "ファイル保存", "自動入力済みファイルはブラウザ経由でダウンロード。サーバーへの保存なし"),
    ]
    lw6 = Inches(2.6)
    rw6 = CW - lw6
    rh6 = Inches(0.68)
    for i, (color, label, val) in enumerate(env_rows):
        ry = CY + Inches(0.08) + rh6 * i
        lb = _rbox(sl, ML, ry, lw6, rh6 - Inches(0.04), color)
        _set_tf(lb, label, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.1), mt=Inches(0.17))
        vb = _rbox(sl, ML + lw6, ry, rw6, rh6 - Inches(0.04),
                   ACCENT2 if i % 2 == 0 else WHITE, GRAY2)
        _box(sl, ML + lw6, ry, Inches(0.06), rh6 - Inches(0.04), color)
        _set_tf(vb, val, size=11, color=DARK, ml=Inches(0.22), mt=Inches(0.17))

    foot6 = _rbox(sl, ML, CY + rh6 * 6 + Inches(0.15), CW, Inches(0.45), NAVY_D)
    _set_tf(foot6,
            "インストール不要・ブラウザだけで完結。既存の PC 環境をそのままご利用いただけます。",
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            ml=Inches(0.2), mt=Inches(0.1))

    # ── Slide 7 : 操作手順 ───────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "5.  操作手順（全体フロー）", 7, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "月初の操作は 7 ステップで完了します（目安：5分以内）", size=10, color=MUTED, italic=True)
    flow7 = [
        ("STEP 1", "アクセス\nログイン",    NAVY,    False),
        ("STEP 2", "対象月\nを確認",        GREEN_D, False),
        ("STEP 3", "勤務時間\n設定",         ORANGE_D,False),
        ("STEP 4", "例外日\n設定（任意）",   ORANGE,  True),
        ("STEP 5", "有給日\n選択（任意）",   PINK,    True),
        ("STEP 6", "Excel\nファイル選択",    TEAL,    False),
        ("STEP 7", "完了・DL",             NAVY,    False),
    ]
    _step_flow(sl, flow7, ML, CY + Inches(0.38), CW, Inches(1.05))
    _tb(sl, ML, CY + Inches(1.52), CW, Inches(0.26),
        "※ STEP 4・5（オレンジ色）は任意。残業・有給がある月のみ操作してください。",
        size=9, color=MUTED, italic=True)
    step_descs7 = [
        (NAVY,    "STEP 1", "ブラウザでアプリのURLを開き、ユーザー名・パスワードを入力してログイン。"),
        (GREEN_D, "STEP 2", "画面上部の「年・月」が自動で当月に設定。別の月は数字を直接変更。"),
        (ORANGE_D,"STEP 3", "月〜金の開始・終了時間を入力。月間スケジュールで全日程・合計稼働時間を確認。"),
        (ORANGE,  "STEP 4", "「＋ 例外日」から残業・休日出勤など個別設定（土日・祝日も上書き可）。"),
        (PINK,    "STEP 5", "「有給を取得する日がある」にチェックしカレンダーで日付を選択。"),
        (TEAL,    "STEP 6", "Excelファイル選択エリアをクリックし、配布された報告書（.xlsx）を選択。"),
        (NAVY,    "STEP 7", "「入力完了・ダウンロード」を押すと自動入力済みファイルを取得。"),
    ]
    dw7 = (CW - Inches(0.3)) / 2
    dy7 = CY + Inches(1.88)
    for i, (color, step, desc) in enumerate(step_descs7):
        col = i % 2
        row = i // 2
        if i == 6:
            dx7, dy7_i, dw7_i = ML, dy7 + Inches(0.95) * 3, CW
        else:
            dx7, dy7_i, dw7_i = ML + col * (dw7 + Inches(0.3)), dy7 + Inches(0.95) * row, dw7
        sb = _rbox(sl, dx7, dy7_i, Inches(0.85), Inches(0.38), color)
        _set_tf(sb, step, size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.07))
        db = _rbox(sl, dx7 + Inches(0.85), dy7_i, dw7_i - Inches(0.85), Inches(0.38),
                   ACCENT2, GRAY2)
        _box(sl, dx7 + Inches(0.85), dy7_i, Inches(0.05), Inches(0.38), color)
        _set_tf(db, desc, size=9, color=DARK, ml=Inches(0.15), mt=Inches(0.07))

    # ── Slide 8 : 自動入力ルール ─────────────────────────
    sl = _slide(prs)
    _chrome(sl, "6.  自動入力ルール", 8, TOTAL)
    _tb(sl, ML, CY + Inches(0.05), CW, Inches(0.28),
        "日の種類に応じて以下のルールで各セルへ自動入力されます", size=10, color=MUTED, italic=True)
    col_w8 = [Inches(2.6), Inches(4.7), Inches(4.5)]
    headers8 = ["日の種類", "開始・終了・休憩時間", "備考欄の記入内容"]
    hx8 = ML
    hy8 = CY + Inches(0.38)
    for w, h in zip(col_w8, headers8):
        cell = _box(sl, hx8, hy8, w, Inches(0.45), NAVY_D)
        _set_tf(cell, h, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.05), mt=Inches(0.1))
        hx8 += w
    rows8 = [
        ("出勤日（平日）",
         "曜日別設定の開始・終了時間\n休憩時間：1時間00分",
         "在宅勤務",           ACCENT2, NAVY_L),
        ("例外日\n（残業・早退・休日出勤）",
         "例外日設定の開始・終了時間\n休憩時間：1時間00分",
         "在宅勤務\n（例外設定が最優先）", ORANGE_L, ORANGE),
        ("有給取得日",
         "空欄",
         "私用により、休暇",   PINK_L, PINK),
        ("祝日",
         "空欄",
         "祝日",               RED_L, RED),
        ("土日",
         "空欄",
         "空欄",               GRAY, MUTED_L),
    ]
    row_h8 = (CH - Inches(0.9)) / len(rows8)
    for ri, (kind, times, note_text, bg, accent_c) in enumerate(rows8):
        ry8 = hy8 + Inches(0.45) + row_h8 * ri
        rx8 = ML
        for ci, (val, w) in enumerate([(kind, col_w8[0]), (times, col_w8[1]), (note_text, col_w8[2])]):
            cell = _rbox(sl, rx8, ry8, w, row_h8, bg, GRAY2)
            if ci == 0:
                _box(sl, rx8, ry8, Inches(0.06), row_h8, accent_c)
            lines8 = val.split('\n')
            _set_tf(cell, lines8, size=11, color=DARK,
                    align=PP_ALIGN.CENTER, ml=Inches(0.12), mt=Inches(0.1))
            rx8 += w
    notes8_y = hy8 + Inches(0.45) + row_h8 * 5 + Inches(0.05)
    for n in [
        "※ 祝日情報は外部API（jpholiday）から自動取得。インターネット未接続時は空欄になる場合があります。",
        "※ 例外日設定は有給・祝日・土日に関わらず最優先で適用されます（休日出勤の記録にも使用可能）。",
    ]:
        _tb(sl, ML, notes8_y, CW, Inches(0.26), n, size=9, color=MUTED, italic=True)
        notes8_y += Inches(0.28)

    # ── Slide 9 : Q&A ────────────────────────────────────
    sl = _slide(prs)
    _chrome(sl, "7.  Q&A（よくある質問）", 9, TOTAL)
    qas9 = [
        ("Q1", "設定は毎月入力し直す必要がありますか？",
         "不要です。曜日別の勤務時間・各種設定はブラウザに自動保存されます。次回起動時には前回の設定が自動的に読み込まれます。"),
        ("Q2", "休日出勤が発生した場合はどうすれば？",
         "「＋ 例外日を追加」から対象日（土日・祝日も可）を選択し、出勤時間を入力します。例外日設定が最優先で反映されます。"),
        ("Q3", "有給を土日・祝日に取得した場合は？",
         "有給取得日のカレンダーで土日・祝日も選択できます。選択した日の備考欄に「私用により、休暇」が記入されます。"),
        ("Q4", "Excelの列構成が違っても使えますか？",
         "はい。アプリがヘッダー文字列を自動検索して対象列を特定するため、列の位置が変わっても正常に動作します。"),
    ]
    qa_h9 = CH / len(qas9)
    for i, (qn_s, q, a) in enumerate(qas9):
        qy9 = CY + qa_h9 * i
        badge9 = _oval(sl, ML, qy9 + qa_h9 * 0.07, Inches(0.55), Inches(0.55), NAVY)
        _set_tf(badge9, qn_s, size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, ml=Inches(0.03), mt=Inches(0.09))
        qrow9 = _rbox(sl, ML + Inches(0.65), qy9 + qa_h9 * 0.06,
                      CW - Inches(0.65), qa_h9 * 0.4, ACCENT2, GRAY2)
        _set_tf(qrow9, q, size=12, bold=True, color=NAVY_D,
                ml=Inches(0.2), mt=Inches(0.09))
        arow9 = _rbox(sl, ML + Inches(0.65), qy9 + qa_h9 * 0.46,
                      CW - Inches(0.65), qa_h9 * 0.5, WHITE, GRAY2)
        _box(sl, ML + Inches(0.65), qy9 + qa_h9 * 0.46, Inches(0.06), qa_h9 * 0.5, NAVY_L)
        _set_tf(arow9, a, size=11, color=DARK, ml=Inches(0.22), mt=Inches(0.1))

    # ── Slide 10 : まとめ・導入のお願い ─────────────────
    sl = _slide(prs)
    _chrome(sl, "8.  まとめ・導入のお願い", 10, TOTAL)
    can_h = _rbox(sl, ML, CY, CW, Inches(0.45), GREEN_D)
    _set_tf(can_h, "本アプリで実現できること", size=13, bold=True, color=WHITE,
            ml=Inches(0.2), mt=Inches(0.1))
    can_items = [
        ("✅", "月初5分以内の操作で1ヶ月分のExcel作業報告書が自動完成",        GREEN_L),
        ("✅", "手入力ゼロ・入力ミスゼロを実現",                               WHITE),
        ("✅", "ブラウザだけで動作。インストール不要・既存PC環境をそのまま利用可能", GREEN_L),
        ("✅", "設定が自動保存されるため、翌月以降は変更のある日だけ修正するだけ",   WHITE),
    ]
    for i, (icon, s, bg) in enumerate(can_items):
        row = _rbox(sl, ML, CY + Inches(0.52) + Inches(0.6) * i, CW, Inches(0.56), bg, GRAY2)
        _box(sl, ML, CY + Inches(0.52) + Inches(0.6) * i, Inches(0.06), Inches(0.56), GREEN)
        _set_tf(row, f"{icon}  {s}", size=12, color=DARK, ml=Inches(0.22), mt=Inches(0.13))

    req_h = _rbox(sl, ML, CY + Inches(3.0), CW, Inches(0.45), NAVY_D)
    _set_tf(req_h, "導入にあたってのお願い", size=13, bold=True, color=WHITE,
            ml=Inches(0.2), mt=Inches(0.1))
    req_body = _rbox(sl, ML, CY + Inches(3.45), CW, Inches(1.55), ACCENT2, ACCENT)
    _set_tf(req_body, [
        "本アプリを社内に導入・展開するにあたり、ご承認をお願いいたします。",
        "現在は主に1名での使用を想定していますが、同一フォーマットを使用する他の社員への展開も容易に行えます。",
        "ご不明な点やご要望がございましたら、お気軽にお申し付けください。",
        "",
        "以上",
    ], size=11, color=DARK, ml=Inches(0.3), mt=Inches(0.18))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════════════
#  開発環境セットアップガイド
#  Google アカウント → GitHub（Google 連携）→ Claude Code（GitHub 連携）→ Codespace 開発
# ════════════════════════════════════════════════════
def _build_setup_guide_pptx():
    prs = _prs()
    TOTAL = 12

    G_BLUE   = RGBColor(0x42, 0x85, 0xF4)   # Google blue
    GH_BLACK = RGBColor(0x24, 0x29, 0x2F)   # GitHub dark
    GH_GRAY  = RGBColor(0xF6, 0xF8, 0xFA)   # GitHub light bg
    ANT_RUST = RGBColor(0xCC, 0x5C, 0x36)   # Anthropic brand
    CS_VIOLT = RGBColor(0x6E, 0x40, 0xC9)   # Codespace purple
    CODE_BG  = RGBColor(0x0D, 0x11, 0x17)   # terminal dark
    CODE_FG  = RGBColor(0x7E, 0xE7, 0x87)   # terminal green
    WARN_BG  = RGBColor(0xFF, 0xFB, 0xEB)

    # ── ローカルヘルパー ──────────────────────────────────────
    def _code_block(slide, x, y, w, h, text):
        bg = _rbox(slide, x, y, w, h, CODE_BG)
        _set_tf(bg, text, size=9, color=CODE_FG, ml=Inches(0.18), mt=Inches(0.1))
        return bg

    def _url_chip(slide, x, y, w, url_text, color=NAVY):
        chip = _rbox(slide, x, y, w, Inches(0.36), ACCENT2, color, line_w=1.0)
        _set_tf(chip, "\U0001f517  " + url_text, size=10, color=color, ml=Inches(0.12), mt=Inches(0.06))
        return chip

    def _step_hdr(slide, num_str, title, color):
        _num_badge(slide, ML, CY + Inches(0.05), Inches(0.58), num_str, color)
        _tb(slide, ML + Inches(0.68), CY + Inches(0.08), CW - Inches(0.68), Inches(0.44),
            title, size=13, bold=True, color=color)

    def _rows(slide, items, start_y, rh=Inches(0.65), lcolor=NAVY_D):
        for i, item in enumerate(items):
            if len(item) == 2:
                icon_s, text = item
                sub = None
            else:
                icon_s, text, sub = item
            bg = ACCENT2 if i % 2 == 0 else WHITE
            h_row = rh + (Inches(0.32) if sub else Inches(0))
            row = _rbox(slide, ML, start_y + rh * i, CW, h_row - Inches(0.04), bg, GRAY2)
            badge = _oval(slide, ML + Inches(0.1), start_y + rh * i + Inches(0.16),
                          Inches(0.3), Inches(0.3), lcolor)
            _set_tf(badge, icon_s, size=7, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, ml=Inches(0.02), mt=Inches(0.04))
            _tb(slide, ML + Inches(0.5), start_y + rh * i + Inches(0.1),
                CW - Inches(0.55), rh - Inches(0.16), text, size=11, color=DARK)
            if sub:
                _code_block(slide, ML + Inches(0.5), start_y + rh * i + rh - Inches(0.08),
                            CW - Inches(0.55), Inches(0.3), sub)

    def _tip(slide, text, color=NAVY_L, bg=None):
        if bg is None:
            bg = ACCENT2
        tip = _rbox(slide, ML, H - FH - Inches(0.52), CW, Inches(0.4), bg, color, line_w=1.0)
        _set_tf(tip, "\U0001f4a1  " + text, size=10, color=DARK, ml=Inches(0.2), mt=Inches(0.08))

    def _warn(slide, text):
        w = _rbox(slide, ML, H - FH - Inches(0.52), CW, Inches(0.4), WARN_BG, ORANGE, line_w=1.0)
        _set_tf(w, "⚠  " + text, size=10, color=ORANGE_D, ml=Inches(0.2), mt=Inches(0.08))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 1: タイトル
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _box(sl, 0, 0, W, H, NAVY_D)
    _oval(sl, W - Inches(5.5), H - Inches(5.5), Inches(7.5), Inches(7.5), RGBColor(0x1E, 0x40, 0xAF))
    _oval(sl, W - Inches(3.8), H - Inches(3.8), Inches(5.5), Inches(5.5), RGBColor(0x1D, 0x4E, 0xD8))
    _oval(sl, -Inches(1.5), -Inches(1.2), Inches(4.5), Inches(4.5), RGBColor(0x1E, 0x40, 0xAF))
    _box(sl, 0, Inches(1.8), Inches(0.18), Inches(2.2), NAVY_L)
    t = _box(sl, Inches(0.38), Inches(1.85), W - Inches(4.8), Inches(1.5), NAVY_D)
    _set_tf(t, "開発環境セットアップガイド", size=30, bold=True, color=WHITE, ml=Inches(0.25), mt=Inches(0.22))
    sub = _box(sl, Inches(0.38), Inches(3.45), W - Inches(4.8), Inches(0.55), NAVY_D)
    _set_tf(sub, "Google アカウント取得からアプリ開発開始まで  ／  所要時間：30〜60 分",
            size=13, color=RGBColor(0xBF, 0xD7, 0xFF), ml=Inches(0.25), mt=Inches(0.06))
    for i, (c, s, tl) in enumerate([
        (G_BLUE,   "STEP 1", "Google アカウントの作成"),
        (GH_BLACK, "STEP 2", "GitHub アカウントの登録（Google アカウント連携）"),
        (ANT_RUST, "STEP 3", "Claude Code の登録（GitHub アカウント連携）"),
        (CS_VIOLT, "STEP 4", "リポジトリ作成と Codespace 起動"),
        (GREEN,    "STEP 5", "Claude Code でアプリ開発開始"),
    ]):
        _oval(sl, Inches(0.45), Inches(4.35) + Inches(0.46) * i, Inches(0.26), Inches(0.26), c)
        _tb(sl, Inches(0.83), Inches(4.32) + Inches(0.46) * i, Inches(6.0), Inches(0.38),
            f"{s}  {tl}", size=10, color=RGBColor(0xBF, 0xD7, 0xFF))
    _box(sl, 0, H - Inches(0.32), W, Inches(0.32), NAVY_D)
    _tb(sl, W - Inches(1.5), H - Inches(0.28), Inches(1.3), Inches(0.25),
        "1  /  12", size=9, color=MUTED_L, align=PP_ALIGN.RIGHT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 2: 全体の流れ
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "全体の流れ", 2, TOTAL)
    _tb(sl, ML, CY + Inches(0.04), CW, Inches(0.28),
        "5 ステップで開発環境を整えます。本ガイドは Codespace（クラウド開発環境）を使った方法の一例です。インストール不要でブラウザだけで完結します。",
        size=10, color=MUTED, italic=True)
    _step_flow(sl, [
        ("STEP 1", "Google\nアカウント作成",          G_BLUE,   False),
        ("STEP 2", "GitHub\n登録\n(Google連携)",       GH_BLACK, False),
        ("STEP 3", "Claude Code\n登録\n(GitHub連携)",  ANT_RUST, False),
        ("STEP 4", "リポジトリ作成\nCodespace起動",   CS_VIOLT, False),
        ("STEP 5", "Claude Code\nで開発開始",          GREEN,    False),
    ], ML, CY + Inches(0.38), CW, Inches(1.2))
    for i, (color, step, head, body) in enumerate([
        (G_BLUE,   "STEP 1", "Google アカウントの作成",
         "Gmail アドレスを取得。STEP 2・3 の登録にもこのアカウントを使います。"),
        (GH_BLACK, "STEP 2", "GitHub アカウントの登録（Google アカウント連携）",
         "Gmail アドレスで GitHub に登録し、設定から Google アカウントを OAuth 連携します。"),
        (ANT_RUST, "STEP 3", "Claude Code の登録（GitHub アカウント連携）",
         "claude.ai に GitHub アカウントで登録。Pro プランで Claude Code がフル活用できます。"),
        (CS_VIOLT, "STEP 4", "リポジトリ作成と Codespace 起動",
         "GitHub にリポジトリを作成し、ブラウザ上の開発環境「Codespace」を起動します。"),
        (GREEN,    "STEP 5", "Claude Code でアプリ開発開始",
         "Codespace のターミナルで Claude Code をインストールし、日本語で指示するだけで開発開始。"),
    ]):
        col = i % 2; row = i // 2
        if i == 4:
            dx, dy, dw = ML, CY + Inches(1.75) + Inches(0.85) * 2, CW
        else:
            dw = (CW - Inches(0.3)) / 2
            dx = ML + col * (dw + Inches(0.3)); dy = CY + Inches(1.75) + Inches(0.85) * row
        s = _rbox(sl, dx, dy, Inches(1.1), Inches(0.38), color)
        _set_tf(s, step, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.05), mt=Inches(0.08))
        h2 = _rbox(sl, dx + Inches(1.1), dy, dw - Inches(1.1), Inches(0.38), ACCENT2, GRAY2)
        _box(sl, dx + Inches(1.1), dy, Inches(0.05), Inches(0.38), color)
        _set_tf(h2, f"{head}  /  {body}", size=9, color=DARK, ml=Inches(0.15), mt=Inches(0.08))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 3: STEP 1 Google アカウント作成
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 1  Google アカウントの作成", 3, TOTAL)
    _step_hdr(sl, "1", "Gmail アドレスを取得します（GitHub・Claude Code の登録にも使います）", G_BLUE)
    _url_chip(sl, ML, CY + Inches(0.58), Inches(5.5), "https://accounts.google.com/signup", G_BLUE)
    _rows(sl, [
        ("①", "上記 URL をブラウザで開き、「アカウントを作成」→「個人用」をクリックします。"),
        ("②", "姓・名・生年月日・性別を順番に入力して「次へ」を押していきます。"),
        ("③", "ユーザー名を決めて入力します（例：yamada.taro2024）。これが Gmail アドレス（@gmail.com）になります。"),
        ("④", "パスワードを設定 → 電話番号で SMS 認証 → 利用規約に同意 → 作成完了です。"),
    ], CY + Inches(1.02), rh=Inches(0.75), lcolor=G_BLUE)
    note3_y = CY + Inches(1.02) + Inches(0.75) * 4 + Inches(0.2)
    nb3 = _rbox(sl, ML, note3_y, CW, Inches(0.72), ACCENT2, G_BLUE, line_w=1.0)
    _set_tf(nb3, [
        "💡  この Gmail アドレスは STEP 2（GitHub）・STEP 3（Claude Code）の登録でも使います。必ずメモしておきましょう。",
        "すでに Google アカウントをお持ちの方は、そのまま STEP 2 へ進んでください。",
    ], size=11, color=NAVY_D, ml=Inches(0.22), mt=Inches(0.12))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 4: STEP 2 GitHub アカウント登録
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 2  GitHub アカウントの登録", 4, TOTAL)
    _step_hdr(sl, "2", "Gmail アドレスで GitHub に登録します", GH_BLACK)
    _url_chip(sl, ML, CY + Inches(0.58), Inches(4.0), "https://github.com/signup", GH_BLACK)
    _rows(sl, [
        ("①", "上記 URL を開き「メールアドレス」欄に STEP 1 の Gmail アドレスを入力して「Continue」。"),
        ("②", "パスワードを設定します（GitHub 専用パスワードを新しく決めてください）。「Continue」を押します。"),
        ("③", "ユーザー名を入力します（例：yamada-taro）。英字・数字・ハイフンのみ使用可。後から変更できます。"),
        ("④", "「Would you like to receive product updates?」→ 「n」で「Continue」。"),
        ("⑤", "「Verify your account」でパズル認証を完了して「Create account」を押します。"),
        ("⑥", "Gmail に届いた件名「Your GitHub launch code」のメールを開き、6桁のコードを入力します。"),
        ("⑦", "「How many team members...」などのアンケートはすべて「Skip personalization」で OK です。"),
    ], CY + Inches(1.02), rh=Inches(0.625), lcolor=GH_BLACK)
    _warn(sl, "ユーザー名は公開されます。本名ではなくニックネームでも構いません。後から変更可能です。")

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 5: STEP 2続 GitHub に Google アカウントを連携 & 2FA
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 2  GitHub 初期設定：Google アカウント連携と 2FA", 5, TOTAL)
    half5 = (CW - Inches(0.3)) / 2
    # 左：Google 連携（簡略・任意）
    gl = _rbox(sl, ML, CY + Inches(0.05), half5, Inches(0.44), G_BLUE)
    _set_tf(gl, "Google アカウントの連携（任意・スキップ可）", size=11, bold=True,
            color=WHITE, ml=Inches(0.15), mt=Inches(0.1))
    for i, (icon_s, text) in enumerate([
        ("①", "右上アイコン → Settings → Password and authentication を開きます。"),
        ("②", "「Social accounts」→「Link Google account」をクリックします。"),
        ("③", "STEP 1 の Google アカウントを選択すれば連携完了です。"),
    ]):
        bg = ACCENT2 if i % 2 == 0 else WHITE
        rr = _rbox(sl, ML, CY + Inches(0.56) + Inches(0.78) * i, half5, Inches(0.74), bg, GRAY2)
        badge = _oval(sl, ML + Inches(0.1), CY + Inches(0.56) + Inches(0.78) * i + Inches(0.22),
                      Inches(0.28), Inches(0.28), G_BLUE)
        _set_tf(badge, icon_s, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.02), mt=Inches(0.04))
        _tb(sl, ML + Inches(0.46), CY + Inches(0.56) + Inches(0.78) * i + Inches(0.16),
            half5 - Inches(0.5), Inches(0.52), text, size=10, color=DARK)
    opt_y = CY + Inches(0.56) + Inches(0.78) * 3 + Inches(0.1)
    opt = _rbox(sl, ML, opt_y, half5, Inches(0.65), ORANGE_L, ORANGE, line_w=0.8)
    _set_tf(opt, ["⚠  この設定は任意です",
                  "スキップしても GitHub・Claude Code の利用に支障はありません。後からいつでも設定できます。"],
            size=9, color=ORANGE_D, ml=Inches(0.15), mt=Inches(0.08))
    _tip(sl, "Google 連携すると「次回から Google でログイン」が使えて便利です（任意）。")
    # 右：2FA
    rx5 = ML + half5 + Inches(0.3)
    gr = _rbox(sl, rx5, CY + Inches(0.05), half5, Inches(0.44), GH_BLACK)
    _set_tf(gr, "2段階認証（2FA）の設定（強く推奨）", size=11, bold=True,
            color=WHITE, ml=Inches(0.15), mt=Inches(0.1))
    for i, (icon_s, text) in enumerate([
        ("①", "「Password and authentication」ページ内「Two-factor authentication」→「Enable」。"),
        ("②", "スマホに「Google Authenticator」アプリをインストールします（App Store / Google Play）。"),
        ("③", "「Authenticator app」を選択し、表示された QR コードをアプリでスキャンします。"),
        ("④", "アプリに表示された 6 桁のコードを GitHub 画面に入力して「Continue」を押します。"),
        ("⑤", "バックアップコードが表示されます。必ずダウンロード・印刷して保管してください。"),
    ]):
        bg = GH_GRAY if i % 2 == 0 else WHITE
        rr = _rbox(sl, rx5, CY + Inches(0.56) + Inches(0.65) * i, half5, Inches(0.61), bg, GRAY2)
        badge = _oval(sl, rx5 + Inches(0.1), CY + Inches(0.56) + Inches(0.65) * i + Inches(0.15),
                      Inches(0.28), Inches(0.28), GH_BLACK)
        _set_tf(badge, icon_s, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.02), mt=Inches(0.04))
        _tb(sl, rx5 + Inches(0.46), CY + Inches(0.56) + Inches(0.65) * i + Inches(0.1),
            half5 - Inches(0.5), Inches(0.48), text, size=10, color=DARK)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 6: STEP 3 claude.ai アカウント作成（GitHub 連携）
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 3  Claude Code の登録（GitHub アカウント連携）", 6, TOTAL)
    _step_hdr(sl, "3", "GitHub アカウントで claude.ai に登録し、Claude Code を有効化します", ANT_RUST)
    _url_chip(sl, ML, CY + Inches(0.58), Inches(3.8), "https://claude.ai", ANT_RUST)
    _rows(sl, [
        ("①", "上記 URL をブラウザで開き、「Sign up」をクリックします。"),
        ("②", "「Continue with GitHub」をクリックします（GitHub アカウントで連携登録）。"),
        ("③", "GitHub のログイン画面が表示されたら、STEP 2 で作成したアカウントでログインします。"),
        ("④", "「Authorize Anthropic」の確認画面が表示されたら「Authorize」をクリックします。"),
        ("⑤", "電話番号認証が求められる場合は、SMS で届いた 6 桁のコードを入力します。"),
        ("⑥", "利用規約に同意して「Continue」→ アカウント作成完了です。"),
    ], CY + Inches(1.02), rh=Inches(0.65), lcolor=ANT_RUST)
    # プランの説明
    plan_y = CY + Inches(1.02) + Inches(0.65) * 6 + Inches(0.12)
    ph = _rbox(sl, ML, plan_y, CW, Inches(0.38), NAVY_D)
    _set_tf(ph, "プランの選択  ─  Claude Code をフル活用するには Pro プランが必要です", size=11, bold=True,
            color=WHITE, ml=Inches(0.2), mt=Inches(0.08))
    pw = (CW - Inches(0.2)) / 2
    for i, (plan, price, desc, bg2, is_rec) in enumerate([
        ("Free プラン", "無料",   "基本的な会話が可能。Claude Code の使用回数に厳しい制限があります。", GRAY, False),
        ("Pro プラン",  "$20/月", "Claude Code 使い放題。本格的なアプリ開発・自動コーディングに対応。",  ACCENT2, True),
    ]):
        px = ML + (pw + Inches(0.2)) * i
        pp = _rbox(sl, px, plan_y + Inches(0.38), pw, Inches(0.4), NAVY if is_rec else GRAY2)
        _set_tf(pp, f"{plan}  /  {price}" + ("  ★推奨" if is_rec else ""),
                size=11, bold=True, color=WHITE if is_rec else DARK, ml=Inches(0.15), mt=Inches(0.08))
        pb = _rbox(sl, px, plan_y + Inches(0.78), pw, Inches(0.42), WHITE, GRAY2)
        _set_tf(pb, desc, size=10, color=DARK, ml=Inches(0.15), mt=Inches(0.1))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 7: STEP 3続 Pro プランへのアップグレード
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 3  Pro プランへのアップグレード", 7, TOTAL)
    _tb(sl, ML, CY + Inches(0.04), CW, Inches(0.3),
        "Claude Code をフル活用するために Pro プラン（$20/月）にアップグレードします",
        size=11, color=MUTED, italic=True)
    up_head = _rbox(sl, ML, CY + Inches(0.38), CW, Inches(0.44), ANT_RUST)
    _set_tf(up_head, "Pro プランへのアップグレード手順", size=13, bold=True,
            color=WHITE, ml=Inches(0.2), mt=Inches(0.09))
    for i, (icon_s, text) in enumerate([
        ("①", "claude.ai にログイン後、左下または画面上部の「Upgrade to Pro」をクリックします。"),
        ("②", "プラン選択画面で「Pro  $20/month」の「Subscribe」ボタンをクリックします。"),
        ("③", "クレジットカード情報を入力して「Subscribe to Pro」をクリックします（Visa / Mastercard 可）。"),
        ("④", "登録完了メールが届き、ダッシュボードに「Pro」バッジが表示されれば完了です。"),
    ]):
        bg = ACCENT2 if i % 2 == 0 else WHITE
        row = _rbox(sl, ML, CY + Inches(0.9) + Inches(0.7) * i, CW, Inches(0.66), bg, GRAY2)
        badge = _oval(sl, ML + Inches(0.1), CY + Inches(0.9) + Inches(0.7) * i + Inches(0.18),
                      Inches(0.3), Inches(0.3), ANT_RUST)
        _set_tf(badge, icon_s, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.02), mt=Inches(0.04))
        _tb(sl, ML + Inches(0.5), CY + Inches(0.9) + Inches(0.7) * i + Inches(0.12),
            CW - Inches(0.55), Inches(0.5), text, size=11, color=DARK)
    # Claude Code 確認
    cc_y = CY + Inches(0.9) + Inches(0.7) * 4 + Inches(0.12)
    cc_head = _rbox(sl, ML, cc_y, CW, Inches(0.44), NAVY_D)
    _set_tf(cc_head, "Claude Code の利用確認（Web 版）", size=13, bold=True,
            color=WHITE, ml=Inches(0.2), mt=Inches(0.09))
    for i, (icon_s, text) in enumerate([
        ("①", "claude.ai にアクセスし、チャット画面が開いていることを確認します。"),
        ("②", "右上のメニューから「Claude Code」または「Use Claude Code」を選択します。"),
        ("③", "コマンド入力欄が表示されれば Pro プランで Claude Code が有効になっています。"),
    ]):
        bg = PURPLE_L if i % 2 == 0 else WHITE
        row = _rbox(sl, ML, cc_y + Inches(0.44) + Inches(0.56) * i, CW, Inches(0.52), bg, GRAY2)
        badge = _oval(sl, ML + Inches(0.1), cc_y + Inches(0.44) + Inches(0.56) * i + Inches(0.12),
                      Inches(0.26), Inches(0.26), PURPLE_D)
        _set_tf(badge, icon_s, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.02), mt=Inches(0.03))
        _tb(sl, ML + Inches(0.44), cc_y + Inches(0.44) + Inches(0.56) * i + Inches(0.08),
            CW - Inches(0.5), Inches(0.42), text, size=11, color=DARK)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 8: STEP 4 リポジトリ作成
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 4  GitHub リポジトリを作成する", 8, TOTAL)
    _step_hdr(sl, "4", "コードを保存・管理する「リポジトリ」を GitHub に作成します", GH_BLACK)
    _url_chip(sl, ML, CY + Inches(0.58), Inches(4.5), "https://github.com/new", GH_BLACK)
    _rows(sl, [
        ("①", "上記 URL を開きます（または GitHub ホーム右上の「＋」→「New repository」をクリック）。"),
        ("②", "「Repository name」に英字でプロジェクト名を入力します（例：my-first-app）。スペース不可。"),
        ("③", "説明（Description）は任意です。スキップして「Public」または「Private」を選択します。"),
        ("④", "「Add a README file」にチェックを入れます（Codespace 起動に必要）。"),
        ("⑤", "「Add .gitignore」→「Python」を選択します（不要ファイルを除外）。"),
        ("⑥", "「Create repository」ボタンをクリックします。"),
    ], CY + Inches(1.02), rh=Inches(0.65), lcolor=GH_BLACK)
    tip8a = _rbox(sl, ML, H - FH - Inches(0.78), CW, Inches(0.25), GH_GRAY, GRAY2)
    _set_tf(tip8a, "Private にしておけば自分だけが閲覧できます。後から Public に変更することも可能です。",
            size=10, color=MUTED, ml=Inches(0.15), mt=Inches(0.04))
    _warn(sl, "README ファイルがないとリポジトリが空の状態になり、Codespace を起動できません。必ずチェックを入れてください。")

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 9: STEP 4続 Codespace 起動
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 4  GitHub Codespace を起動する", 9, TOTAL)
    _tb(sl, ML, CY + Inches(0.04), CW, Inches(0.3),
        "本ガイドでは Codespace（ブラウザ上の VS Code 環境）を使う方法を紹介します。Codespace はあくまで一例であり、VS Code のローカルインストールなど他の方法でも問題ありません。",
        size=11, color=MUTED, italic=True)
    # What is Codespace
    cs_info_y = CY + Inches(0.38)
    ci = _rbox(sl, ML, cs_info_y, CW, Inches(0.44), CS_VIOLT)
    _set_tf(ci, "Codespace を使うメリット", size=12, bold=True, color=WHITE, ml=Inches(0.2), mt=Inches(0.09))
    feat_w = (CW - Inches(0.3)) / 3
    for i, (icon_s, head, body) in enumerate([
        ("\U0001f4bb", "ブラウザだけで動く", "VS Code がそのままブラウザで使えます。インストール不要。"),
        ("\U000026a1", "環境が最初から整っている", "Node.js・Python・Git がプリインストール済みです。"),
        ("\U00002601", "クラウド保存",         "作業内容はクラウドに保存。どのPCからでもアクセス可。"),
    ]):
        fx = ML + (feat_w + Inches(0.15)) * i
        fb = _rbox(sl, fx, cs_info_y + Inches(0.44), feat_w, Inches(0.9), ACCENT2, GRAY2)
        _tb(sl, fx + Inches(0.05), cs_info_y + Inches(0.5), feat_w - Inches(0.1), Inches(0.3),
            icon_s + "  " + head, size=11, bold=True, color=CS_VIOLT)
        _tb(sl, fx + Inches(0.1), cs_info_y + Inches(0.82), feat_w - Inches(0.15), Inches(0.5),
            body, size=10, color=DARK)
    # 起動手順
    launch_y = cs_info_y + Inches(1.42)
    lh = _rbox(sl, ML, launch_y, CW, Inches(0.44), GH_BLACK)
    _set_tf(lh, "Codespace の起動手順", size=12, bold=True, color=WHITE, ml=Inches(0.2), mt=Inches(0.09))
    for i, (icon_s, text) in enumerate([
        ("①", "作成したリポジトリのページを開きます（github.com → Your repositories → リポジトリ名）。"),
        ("②", "緑色の「<> Code」ボタンをクリックし、「Codespaces」タブを選択します。"),
        ("③", "「Create codespace on main」をクリックします。"),
        ("④", "ブラウザ内に VS Code が起動します（初回は 1〜2 分かかります）。"),
        ("⑤", "画面下部の「Terminal」タブをクリックしてターミナルを開きます。表示されない場合は メニュー「Terminal」→「New Terminal」。"),
    ]):
        bg = ACCENT2 if i % 2 == 0 else WHITE
        rr = _rbox(sl, ML, launch_y + Inches(0.44) + Inches(0.62) * i, CW, Inches(0.58), bg, GRAY2)
        badge = _oval(sl, ML + Inches(0.1), launch_y + Inches(0.44) + Inches(0.62) * i + Inches(0.14),
                      Inches(0.28), Inches(0.28), GH_BLACK)
        _set_tf(badge, icon_s, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.02), mt=Inches(0.04))
        _tb(sl, ML + Inches(0.46), launch_y + Inches(0.44) + Inches(0.62) * i + Inches(0.08),
            CW - Inches(0.5), Inches(0.46), text, size=10, color=DARK)
    _tip(sl, "無料プランでは月 120 時間・ストレージ 15 GB が使用できます。Pro プランは月 180 時間まで無料。")

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 10: STEP 4続 Claude Code インストール & 認証
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 4  Codespace 内で Claude Code をインストールする", 10, TOTAL)
    _tb(sl, ML, CY + Inches(0.04), CW, Inches(0.28),
        "Codespace のターミナルに以下のコマンドを順番に入力してください",
        size=11, color=MUTED, italic=True)
    cmd_y = CY + Inches(0.38)
    half10 = (CW - Inches(0.3)) / 2
    # 左列：コマンド手順
    for i, (head, cmd, note) in enumerate([
        ("① Node.js のバージョンを確認する",
         "node --version\n# 出力例: v20.x.x または v22.x.x",
         "v18 以上であれば問題ありません。Codespace には最初から入っています。"),
        ("② Claude Code をインストールする",
         "npm install -g @anthropic-ai/claude-code",
         "「claude」コマンドが使えるようになります。数十秒かかります。"),
        ("③ インストールを確認する",
         "claude --version\n# 出力例: 1.x.x",
         "バージョン番号が表示されればインストール成功です。"),
        ("④ 初回認証を行う",
         "claude\n# → ブラウザに認証URLが表示されます",
         "表示された URL をクリックし、claude.ai の Pro アカウントでログインします。"),
    ]):
        hy = cmd_y + Inches(1.4) * i
        hb = _rbox(sl, ML, hy, half10, Inches(0.35), GH_BLACK)
        _set_tf(hb, head, size=10, bold=True, color=WHITE, ml=Inches(0.12), mt=Inches(0.07))
        _code_block(sl, ML, hy + Inches(0.35), half10, Inches(0.6), cmd)
        nb = _rbox(sl, ML, hy + Inches(0.95), half10, Inches(0.38), ACCENT2, GRAY2)
        _set_tf(nb, note, size=9, color=DARK, ml=Inches(0.12), mt=Inches(0.07))
    # 右列：認証フロー説明（詳細）
    rx10 = ML + half10 + Inches(0.3)
    ah = _rbox(sl, rx10, cmd_y, half10, Inches(0.44), ANT_RUST)
    _set_tf(ah, "初回認証の流れ（詳細）", size=12, bold=True, color=WHITE, ml=Inches(0.18), mt=Inches(0.09))
    for i, (icon_s, text) in enumerate([
        ("\U0001f4bb", "ターミナルで「claude」と入力して Enter を押します。"),
        ("\U0001f517", "「Please open the following URL in your browser:」に続き、長い URL が表示されます。"),
        ("\U0001f5b1", "Ctrl キーを押しながら URL をクリック（またはURLをコピーしてブラウザのアドレスバーに貼り付け）します。"),
        ("\U0001f511", "ブラウザに claude.ai のログイン画面が表示されます。Pro アカウントでログインします。"),
        ("🔐",     "「Claude Code wants to access your account」という画面で「Authorize」をクリックします。"),
        ("✅",     "「Authentication successful! You may close this tab.」と表示されたらタブを閉じます。"),
        ("\U0001f4ac", "ターミナルに「✓ Logged in as 〇〇」と表示されます。"),
        ("▶",      "「>」プロンプトが出たら準備完了！そのまま日本語で指示を入力してください。"),
    ]):
        bg = ACCENT2 if i % 2 == 0 else WHITE
        ar = _rbox(sl, rx10, cmd_y + Inches(0.44) + Inches(0.47) * i, half10, Inches(0.43), bg, GRAY2)
        _tb(sl, rx10 + Inches(0.1), cmd_y + Inches(0.44) + Inches(0.47) * i + Inches(0.04),
            Inches(0.3), Inches(0.32), icon_s, size=13, align=PP_ALIGN.CENTER)
        _tb(sl, rx10 + Inches(0.46), cmd_y + Inches(0.44) + Inches(0.47) * i + Inches(0.06),
            half10 - Inches(0.5), Inches(0.34), text, size=9, color=DARK)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 11: STEP 5 開発開始とコード保存
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "STEP 5  Claude Code でアプリ開発を開始する", 11, TOTAL)
    _tb(sl, ML, CY + Inches(0.04), CW, Inches(0.3),
        "あとは Claude Code に日本語で指示するだけです！（Codespace のターミナルで実行）",
        size=13, bold=True, color=GREEN_D)
    half11 = (CW - Inches(0.3)) / 2
    # 左：指示の例
    eh = _rbox(sl, ML, CY + Inches(0.42), half11, Inches(0.44), GREEN_D)
    _set_tf(eh, "Claude Code への指示の例", size=12, bold=True, color=WHITE,
            ml=Inches(0.2), mt=Inches(0.09))
    for i, ex in enumerate([
        "「Flask を使ったシンプルな Todo アプリを作ってください」",
        "「このコードにコメントを日本語で追加してください」",
        "「バグを見つけて直してください」",
        "「README.md を日本語で書いてください」",
        "「ログイン機能を追加してください」",
    ]):
        bg = GREEN_L if i % 2 == 0 else WHITE
        er = _rbox(sl, ML, CY + Inches(0.94) + Inches(0.55) * i, half11, Inches(0.51), bg, GRAY2)
        _box(sl, ML, CY + Inches(0.94) + Inches(0.55) * i, Inches(0.05), Inches(0.51), GREEN)
        _set_tf(er, "\U0001f4ac  " + ex, size=10, color=DARK, ml=Inches(0.2), mt=Inches(0.12))
    # 右：git コマンド
    rx11 = ML + half11 + Inches(0.3)
    gh = _rbox(sl, rx11, CY + Inches(0.42), half11, Inches(0.44), GH_BLACK)
    _set_tf(gh, "作業後は GitHub に保存（push）", size=12, bold=True, color=WHITE,
            ml=Inches(0.18), mt=Inches(0.09))
    for i, (cmd, note) in enumerate([
        ("git status",                          "変更されたファイルを確認します。"),
        ('git add .',                           "すべての変更をステージングします。"),
        ('git commit -m "変更内容の説明"',      "変更をコミット（記録）します。"),
        ("git push",                            "GitHub に変更をアップロードします。"),
    ]):
        bg = ACCENT2 if i % 2 == 0 else WHITE
        cr = _rbox(sl, rx11, CY + Inches(0.94) + Inches(0.65) * i, half11, Inches(0.61), bg, GRAY2)
        _code_block(sl, rx11 + Inches(0.1), CY + Inches(0.94) + Inches(0.65) * i + Inches(0.06),
                    half11 - Inches(0.2), Inches(0.3), cmd)
        _tb(sl, rx11 + Inches(0.1), CY + Inches(0.94) + Inches(0.65) * i + Inches(0.36),
            half11 - Inches(0.15), Inches(0.22), note, size=9, color=MUTED)
    _tip(sl, "Codespace 内の変更は自動保存されます。git push するまでは GitHub 側には反映されません。", NAVY_L, ACCENT2)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    #  Slide 12: Q&A
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sl = _slide(prs)
    _chrome(sl, "よくある質問（Q&A）", 12, TOTAL)
    qa_h = CH / 4
    for i, (qn_s, q, a, color) in enumerate([
        ("Q1", "Codespace は無料で使えますか？",
         "個人アカウントは月120時間・15GBまで無料です。Pro プランユーザーは月180時間まで無料。超過後は課金が発生しますが、通常の使用では無料枠内に収まります。",
         G_BLUE),
        ("Q2", "Claude Code の Pro プランは必須ですか？",
         "Free プランでも試用できますが、使用回数に厳しい制限があります。本格的な開発には Pro プラン（$20/月）を推奨します。Claude Code で自動生成できるコード量が大幅に増えます。",
         ANT_RUST),
        ("Q3", "Codespace を閉じてもデータは消えませんか？",
         "ブラウザを閉じても Codespace は保持されます。GitHub の「Codespaces」ページから再開できます。ただし非アクティブ状態が 30 日続くと自動削除されます。",
         CS_VIOLT),
        ("Q4", "GitHub・claude.ai のパスワードを忘れた場合は？",
         "GitHub：ログイン画面の「Forgot password?」からメール（Gmail）でリセットできます。claude.ai：「Forgot password?」または「Continue with GitHub」を使って再ログインできます。",
         GH_BLACK),
    ]):
        qy = CY + qa_h * i
        badge = _oval(sl, ML, qy + qa_h * 0.1, Inches(0.55), Inches(0.55), color)
        _set_tf(badge, qn_s, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                ml=Inches(0.03), mt=Inches(0.09))
        qrow = _rbox(sl, ML + Inches(0.65), qy + qa_h * 0.08,
                     CW - Inches(0.65), qa_h * 0.38, ACCENT2, GRAY2)
        _box(sl, ML + Inches(0.65), qy + qa_h * 0.08, Inches(0.06), qa_h * 0.38, color)
        _set_tf(qrow, q, size=11, bold=True, color=NAVY_D, ml=Inches(0.22), mt=Inches(0.08))
        arow = _rbox(sl, ML + Inches(0.65), qy + qa_h * 0.46,
                     CW - Inches(0.65), qa_h * 0.5, WHITE, GRAY2)
        _box(sl, ML + Inches(0.65), qy + qa_h * 0.46, Inches(0.06), qa_h * 0.5, MUTED_L)
        _set_tf(arow, a, size=10, color=DARK, ml=Inches(0.22), mt=Inches(0.07))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
